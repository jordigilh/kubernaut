#!/usr/bin/env python3
"""Update kubernaut-technical-overview.pptx from the original Google Slides export.

Strategy: Load the original PPTX (preserving its theme/template exactly),
apply text fixes to existing slides, and insert new slides.

Usage:
    /tmp/pptx-venv/bin/python scripts/generate_slides.py

Input:
    ~/Downloads/Kubernaut Technical Overview.pptx  (Google Slides export)

Output:
    docs/presentations/kubernaut-technical-overview.pptx
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SCRIPT_DIR = Path(__file__).resolve().parent
PROJECT_DIR = SCRIPT_DIR.parent
TEMPLATE = Path.home() / "Downloads" / "Kubernaut Technical Overview.pptx"
OUTPUT = PROJECT_DIR / "docs" / "presentations" / "kubernaut-technical-overview.pptx"

# ---------------------------------------------------------------------------
# Theme colours (for new elements only)
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x2B, 0x57, 0x9A)
TABLE_HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
TABLE_ALT_ROW = RGBColor(0xF0, 0xF4, 0xF8)

FONT_BODY = "Calibri"
FONT_CODE = "Courier New"

# Original slide positions (matching the title text box + accent bar)
TITLE_LEFT = Emu(548655)
TITLE_TOP = Emu(274320)
TITLE_WIDTH = Emu(7544100)
TITLE_HEIGHT = Emu(480000)
BAR_LEFT = Emu(548655)
BAR_TOP = Emu(754380)
BAR_WIDTH = Emu(1371600)
BAR_HEIGHT = Emu(28500)
CONTENT_LEFT = Emu(548655)
CONTENT_TOP = Emu(960120)


# ---------------------------------------------------------------------------
# Text replacement helpers
# ---------------------------------------------------------------------------

def _find_and_replace_in_shape(shape, old_text, new_text):
    """Replace text in a shape's text frame, preserving formatting."""
    if not shape.has_text_frame:
        return False
    found = False
    for paragraph in shape.text_frame.paragraphs:
        full_text = paragraph.text
        if old_text in full_text:
            new_full = full_text.replace(old_text, new_text)
            if paragraph.runs:
                paragraph.runs[0].text = new_full
                for run in paragraph.runs[1:]:
                    run.text = ""
            found = True
    return found


def _replace_paragraph_text(paragraph, new_text):
    """Replace all text in a paragraph, preserving its first run's formatting."""
    if paragraph.runs:
        paragraph.runs[0].text = new_text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = new_text


def _get_code_shape(slide):
    """Find the code block shape (AUTO_SHAPE with text) on a slide."""
    for shape in slide.shapes:
        if shape.shape_type == 1 and shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and len(text) > 50:
                return shape
    return None


def _rewrite_code_block(shape, new_code):
    """Replace the entire code content in a code block shape."""
    tf = shape.text_frame
    lines = new_code.strip().split("\n")

    existing_paras = list(tf.paragraphs)

    for i, line in enumerate(lines):
        if i < len(existing_paras):
            _replace_paragraph_text(existing_paras[i], line)
        else:
            p = tf.add_paragraph()
            p.text = line
            if existing_paras:
                ref = existing_paras[0]
                if ref.runs:
                    for run in p.runs:
                        run.font.size = ref.runs[0].font.size
                        run.font.name = ref.runs[0].font.name

    for i in range(len(lines), len(existing_paras)):
        _replace_paragraph_text(existing_paras[i], "")


# ---------------------------------------------------------------------------
# Slide insertion helpers
# ---------------------------------------------------------------------------

def _duplicate_slide(prs, source_index):
    """Duplicate a slide (deep copy) and append it at the end.

    Returns the new slide object.
    """
    source = prs.slides[source_index]
    slide_layout = source.slide_layout

    new_slide = prs.slides.add_slide(slide_layout)

    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    for shape in source.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    bg_src = source._element.find(qn("p:bg"))
    if bg_src is not None:
        bg_new = copy.deepcopy(bg_src)
        existing_bg = new_slide._element.find(qn("p:bg"))
        if existing_bg is not None:
            new_slide._element.replace(existing_bg, bg_new)
        else:
            new_slide._element.insert(0, bg_new)

    return new_slide


def _move_slide_to(prs, from_index, to_index):
    """Move a slide from one position to another in the slide list."""
    sldIdLst = prs.slides._sldIdLst
    entries = list(sldIdLst)
    entry = entries[from_index]
    sldIdLst.remove(entry)
    entries_after = list(sldIdLst)
    if to_index >= len(entries_after):
        sldIdLst.append(entry)
    else:
        entries_after[to_index].addprevious(entry)


def _clear_slide_shapes(slide):
    """Remove all shapes from a slide."""
    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)


def _add_accent_bar(slide):
    """Add the accent bar matching the original template."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, BAR_LEFT, BAR_TOP, BAR_WIDTH, BAR_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    shape.shadow.inherit = False


def _add_title_box(slide, line1, line2=None):
    """Add a title matching the original's format."""
    txBox = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = line1
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = FONT_BODY
    if line2:
        p2 = tf.add_paragraph()
        p2.text = line2
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.font.name = FONT_BODY
    _add_accent_bar(slide)


def _add_textbox(slide, text, left, top, width, height,
                 font_size=10, bold=False, color=DARK_GRAY, font_name=None):
    """Add a generic text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name or FONT_BODY
    return txBox


def _add_bullets(slide, items, left, top, width, height, font_size=9):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.font.name = FONT_BODY
        p.space_after = Pt(3)
        pPr = p._p.get_or_add_pPr()
        for existing in pPr.findall(qn("a:buChar")):
            pPr.remove(existing)
        for existing in pPr.findall(qn("a:buNone")):
            pPr.remove(existing)
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "\u2022"})
        pPr.append(buChar)
    return txBox


def _add_code_block(slide, code_text, left, top, width, height, font_size=7):
    """Add a code block with background rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()
    shape.shadow.inherit = False

    txBox = slide.shapes.add_textbox(
        left + Emu(80000), top + Emu(50000),
        width - Emu(160000), height - Emu(100000),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = FONT_CODE
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(1)
        p.space_before = Pt(0)
    return txBox


# ---------------------------------------------------------------------------
# Fix existing slides
# ---------------------------------------------------------------------------

def fix_signal_processing(slide):
    """Slide 3: Fix rego.v2 -> v1, already has development env_rank and P2."""
    for shape in slide.shapes:
        _find_and_replace_in_shape(shape, "import rego.v2", "import rego.v1")
        _find_and_replace_in_shape(shape,
            "Recovery context propagation for retries",
            "Degraded mode: partial enrichment when target not found")


def fix_ai_analysis(slide):
    """Slide 4: Fix rego.v2 -> v1, update approval rule."""
    for shape in slide.shapes:
        _find_and_replace_in_shape(shape, "import rego.v2", "import rego.v1")
        _find_and_replace_in_shape(shape,
            "Production + unvalidated target requires approval",
            "Missing affected resource: default-deny (ADR-055)")
        _find_and_replace_in_shape(shape, "not target_validated", "not has_affected_resource")

    code_shape = _get_code_shape(slide)
    if code_shape:
        tf = code_shape.text_frame
        for para in tf.paragraphs:
            text = para.text
            if "not target_validated" in text:
                _replace_paragraph_text(para, "    not has_affected_resource")
            if "Production + unvalidated target" in text:
                _replace_paragraph_text(para,
                    "# Missing affected resource: default-deny (ADR-055)")
            # Remove the is_production line that was paired with target_validated
            # This is trickier - we need to find the specific is_production
            # that comes right before not target_validated in the approval rule

        # Find and remove the extra is_production line in the 3rd rule
        paras = list(tf.paragraphs)
        for i, para in enumerate(paras):
            if "not has_affected_resource" in para.text:
                # The line before should be is_production - remove it
                if i > 0 and "is_production" in paras[i-1].text:
                    # Check it's the right one (in the 3rd require_approval block)
                    _replace_paragraph_text(paras[i-1], "")


def fix_hapi(slide):
    """Slide 5: Add forward references to new slides."""
    for shape in slide.shapes:
        _find_and_replace_in_shape(shape,
            "list_available_actions",
            "list_available_actions (see Slide 7 for taxonomy)")
        _find_and_replace_in_shape(shape,
            "fetch schema and parameters",
            "fetch schema and parameters (see Slide 10)")


def fix_workflow_execution(slide):
    """Slide 7 (original), will become slide 8 after insert."""
    pass  # Already has "Pure executor" line in the original


# ---------------------------------------------------------------------------
# Build new slides
# ---------------------------------------------------------------------------

def build_action_type_taxonomy(prs, donor_index):
    """Create the Action Type Taxonomy slide by cloning a donor and replacing content."""
    new_slide = _duplicate_slide(prs, donor_index)
    _clear_slide_shapes(new_slide)

    _add_title_box(new_slide,
        "Action Type Taxonomy",
        "V1.0 (DD-WORKFLOW-016)  \u2022  10 curated remediation action types")

    rows = [
        ("Action Type", "What It Does"),
        ("ScaleReplicas", "Horizontally scale a workload by adjusting the replica count"),
        ("RestartPod", "Kill and recreate one or more pods"),
        ("IncreaseCPULimits", "Increase CPU resource limits on containers"),
        ("IncreaseMemoryLimits", "Increase memory resource limits on containers"),
        ("RollbackDeployment", "Revert a deployment to its previous stable revision"),
        ("DrainNode", "Drain and cordon a node, evicting all pods"),
        ("CordonNode", "Cordon a node to prevent new scheduling without eviction"),
        ("RestartDeployment", "Rolling restart of all pods in a workload"),
        ("CleanupNode", "Reclaim disk space by purging temp files and unused images"),
        ("DeletePod", "Delete pods stuck in a terminal state"),
    ]

    n_rows = len(rows)
    table_shape = new_slide.shapes.add_table(
        n_rows, 2, CONTENT_LEFT, CONTENT_TOP, Emu(4800000), Emu(3600000))
    table = table_shape.table
    table.columns[0].width = Emu(1700000)
    table.columns[1].width = Emu(3100000)

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(8)
            p.font.name = FONT_BODY
            if row_idx == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
                p.alignment = PP_ALIGN.CENTER
            else:
                p.font.color.rgb = DARK_GRAY
                if col_idx == 0:
                    p.font.bold = True
                    p.font.name = FONT_CODE
                    p.font.size = Pt(7)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = TABLE_ALT_ROW
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    _add_textbox(new_slide, "Three-Step Discovery Example (OOMKill scenario)",
                 Emu(5400000), CONTENT_TOP, Emu(3600000), Emu(200000),
                 font_size=9, bold=True, color=ACCENT_BLUE)

    discovery_code = (
        'Step 1: list_available_actions(\n'
        '          severity=critical,\n'
        '          environment=staging)\n'
        '  LLM: "What actions are available?"\n'
        '  -> IncreaseMemoryLimits, ScaleReplicas,\n'
        '     RestartPod, RollbackDeployment, ...\n'
        '\n'
        'Step 2: list_workflows(\n'
        '          action_type=IncreaseMemoryLimits,\n'
        '          severity=critical)\n'
        '  LLM: "Which workflows implement this?"\n'
        '  -> oomkill-increase-memory-v1 (job)\n'
        '     oomkill-increase-memory-v1 (tekton)\n'
        '\n'
        'Step 3: get_workflow(\n'
        '          workflow_id=oomkill-increase-memory-v1,\n'
        '          version=1.0.0)\n'
        '  LLM: "Give me the schema."\n'
        '  -> full workflow-schema.yaml (Slide 10)\n'
        '\n'
        'LLM populates parameters:\n'
        '  MEMORY_LIMIT_NEW=128Mi'
    )
    _add_code_block(new_slide, discovery_code,
                    Emu(5400000), Emu(1200000), Emu(3600000), Emu(3400000),
                    font_size=7)

    return new_slide


def build_workflow_schema(prs, donor_index):
    """Create the Workflow Schema slide by cloning a donor and replacing content."""
    new_slide = _duplicate_slide(prs, donor_index)
    _clear_slide_shapes(new_slide)

    _add_title_box(new_slide,
        "Workflow Schema",
        "workflow-schema.yaml  \u2022  OCI image  \u2022  tag + digest (DD-WORKFLOW-002 v2.4)")

    schema_yaml = (
        'metadata:\n'
        '  workflowId: oomkill-increase-memory-v1\n'
        '  version: "1.0.0"\n'
        '  description:\n'
        '    what: "Increases memory limits for OOMKill pods"\n'
        '    whenToUse: "Pods OOMKilled due to low limits"\n'
        '    whenNotToUse: "OOM caused by a memory leak"\n'
        '    preconditions: "Managed by Deployment or SS"\n'
        '\n'
        'actionType: IncreaseMemoryLimits\n'
        '\n'
        'labels:\n'
        '  signalType: OOMKilled\n'
        '  severity: [critical, high]\n'
        '  environment: [production, staging, test]\n'
        '  component: "*"\n'
        '  priority: "*"\n'
        '\n'
        'execution:\n'
        '  engine: job\n'
        '  # V1.0: tag + SHA256 digest for audit trail\n'
        '  containerImage: quay.io/kubernaut-cicd/\n'
        '    workflows/oomkill-increase-memory:v1.0.0\n'
        '    @sha256:a3ed95caeb02ffe68cdd9fd8440...\n'
        '\n'
        'customLabels:\n'
        '  team: platform-sre\n'
        '  cost-profile: memory-intensive\n'
        '  change-risk: low\n'
        '\n'
        'parameters:\n'
        '  - name: TARGET_RESOURCE_KIND\n'
        '    type: string, required: true\n'
        '  - name: TARGET_RESOURCE_NAME\n'
        '    type: string, required: true\n'
        '  - name: TARGET_NAMESPACE\n'
        '    type: string, required: true\n'
        '  - name: MEMORY_LIMIT_NEW\n'
        '    type: string, required: true\n'
        '    description: "New memory limit (128Mi, 1Gi)"'
    )
    _add_code_block(new_slide, schema_yaml,
                    CONTENT_LEFT, CONTENT_TOP, Emu(4800000), Emu(3900000),
                    font_size=7)

    _add_textbox(new_slide, "Key Design Decisions",
                 Emu(5700000), CONTENT_TOP, Emu(3400000), Emu(200000),
                 font_size=10, bold=True, color=ACCENT_BLUE)

    _add_bullets(new_slide, [
        "containerImage with tag + digest ensures exact workflow binary is traceable in the audit trail",
        "actionType maps to the action taxonomy used by the LLM's three-step discovery protocol",
        "labels (mandatory) enable catalog filtering: severity, environment, component, priority",
        "customLabels (optional) are operator-defined key-value pairs for team/org-specific filtering \u2014 stored in their own JSONB column",
        "description fields (what, whenToUse, whenNotToUse, preconditions) are shown to the LLM during workflow selection",
        "parameters are populated by the LLM based on the incident context",
    ],
        Emu(5700000), Emu(1200000), Emu(3400000), Emu(3700000),
        font_size=8,
    )

    return new_slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    if not TEMPLATE.exists():
        print(f"ERROR: Template not found at {TEMPLATE}")
        print("Download from Google Slides first.")
        return

    print(f"Loading template: {TEMPLATE}")
    prs = Presentation(str(TEMPLATE))
    n_original = len(prs.slides)
    print(f"Original slides: {n_original}")

    # ---------------------------------------------------------------
    # Step 1: Fix text in existing slides (in-place, preserving format)
    # ---------------------------------------------------------------
    print("Fixing Signal Processing (slide 3): rego.v2 -> v1")
    fix_signal_processing(prs.slides[3])

    print("Fixing AI Analysis (slide 4): rego.v2 -> v1, approval rule")
    fix_ai_analysis(prs.slides[4])

    print("Fixing HAPI (slide 5): forward references")
    fix_hapi(prs.slides[5])

    # ---------------------------------------------------------------
    # Step 2: Create new slides (appended at end, then moved)
    # ---------------------------------------------------------------
    # Use slide 7 (Workflow Execution) as donor for layout cloning
    # since it's a simple bullet-list slide with the right template.

    print("Creating Action Type Taxonomy slide")
    build_action_type_taxonomy(prs, donor_index=7)
    # Now at index 15 (end), needs to go to position 6 (after HAPI)

    print("Creating Workflow Schema slide")
    build_workflow_schema(prs, donor_index=7)
    # Now at index 16 (end), needs to go to position 9 (after WFE)

    # ---------------------------------------------------------------
    # Step 3: Reorder slides
    # ---------------------------------------------------------------
    # Current order (0-indexed):
    #   0:Title, 1:Pipeline, 2:Gateway, 3:SP, 4:AA, 5:HAPI,
    #   6:RO, 7:WFE, 8:EM, 9:Notif, 10:DS, 11:Auth, 12:Cov,
    #   13:Demo, 14:Q&A, 15:ActionTaxonomy(new), 16:WorkflowSchema(new)
    #
    # Target order:
    #   0:Title, 1:Pipeline, 2:Gateway, 3:SP, 4:AA, 5:HAPI,
    #   6:ActionTaxonomy, 7:RO, 8:WFE, 9:WorkflowSchema,
    #   10:EM, 11:Notif, 12:DS, 13:Auth, 14:Cov, 15:Demo, 16:Q&A

    print("Reordering slides...")
    # Move ActionTaxonomy (currently 15) to position 6
    _move_slide_to(prs, 15, 6)
    # After that move: WorkflowSchema is at 16, target pos is 9
    # But ActionTaxonomy pushed everything after 6 down by 1, so
    # WFE is now at 8 and we want WorkflowSchema right after it (pos 9)
    _move_slide_to(prs, 16, 9)

    # ---------------------------------------------------------------
    # Verify
    # ---------------------------------------------------------------
    print(f"\nFinal slide order ({len(prs.slides)} slides):")
    for i, slide in enumerate(prs.slides):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text[:50].replace("\n", " | ")
                if t.strip():
                    title = t
                    break
        print(f"  {i:2d}: {title}")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"\nGenerated: {OUTPUT}")
    print(f"Size:      {OUTPUT.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    main()
